"""
RAG Engine — uses FAISS instead of ChromaDB (more stable).
Supports: PDF, DOCX, TXT, MD, Excel, CSV, PPTX, HTML, JSON

Arabic-fixed version:
  - PDF text extracted with PyMuPDF (fitz) for correct RTL/word order.
  - Arabic presentation-forms reshaping/bidi reversal repaired on load.
  - Multilingual embedding model (Arabic + English).
"""
from __future__ import annotations
import json, logging, os, pickle, re, unicodedata
from pathlib import Path
from typing import Any
import httpx
import numpy as np

logger = logging.getLogger(__name__)

GROQ_URL    = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL  = "llama-3.3-70b-versatile"
GROQ_KEY    = os.getenv("GROQ_API_KEY", "")
# Multilingual model: strong on Arabic AND English.
EMBED_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls",
    ".csv", ".pptx", ".ppt", ".html", ".htm",
    ".json", ".txt", ".md"
}

_embedder = None

def _get_embedder():
    global _embedder
    if _embedder is None:
        from sentence_transformers import SentenceTransformer
        _embedder = SentenceTransformer(EMBED_MODEL)
    return _embedder


# ── Arabic Text Repair ────────────────────────────────────────────────────────

# Arabic presentation forms (FB50–FEFF / FE70–FEFF) appear when a PDF stores
# glyphs instead of logical characters. NFKC maps them back to base letters.
_AR_PRESENTATION = re.compile(r"[\uFB50-\uFDFF\uFE70-\uFEFF]")
_AR_RANGE        = re.compile(r"[\u0600-\u06FF\u0750-\u077F\uFB50-\uFDFF\uFE70-\uFEFF]")


def _has_arabic(text: str) -> bool:
    return bool(_AR_RANGE.search(text))


def _looks_reversed(text: str) -> bool:
    """
    Heuristic: presentation-form glyphs in the raw text usually mean the PDF
    stored shaped+visually-ordered Arabic, which extracts reversed.
    """
    return bool(_AR_PRESENTATION.search(text))


def _fix_arabic(text: str) -> str:
    """
    Repair Arabic extracted from PDFs:
      1) Normalize presentation forms -> base Arabic letters (NFKC).
      2) If the line was stored visually (reversed), restore logical order
         using arabic_reshaper + python-bidi when available; otherwise fall
         back to a per-line word/char reversal heuristic.
    Non-Arabic text is returned untouched.
    """
    if not text or not _has_arabic(text):
        return text

    needs_bidi = _looks_reversed(text)

    # Step 1: normalize presentation forms to base letters.
    text = unicodedata.normalize("NFKC", text)

    if not needs_bidi:
        return text

    # Step 2: try the proper libraries first.
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display

        fixed_lines = []
        for line in text.split("\n"):
            if _has_arabic(line):
                # get_display converts visual->logical (and vice-versa);
                # reshaping keeps letter forms consistent.
                reshaped = arabic_reshaper.reshape(line)
                fixed_lines.append(get_display(reshaped))
            else:
                fixed_lines.append(line)
        return "\n".join(fixed_lines)
    except Exception as e:
        logger.warning("arabic_reshaper/bidi unavailable (%s); using fallback", e)

    # Step 3: fallback — reverse the visual order of each Arabic line so that
    # logical reading order is restored. Latin tokens inside are kept readable.
    fixed_lines = []
    for line in text.split("\n"):
        if _has_arabic(line):
            tokens = line.split(" ")
            fixed_lines.append(" ".join(reversed(tokens)))
        else:
            fixed_lines.append(line)
    return "\n".join(fixed_lines)


# ── Document Loaders ─────────────────────────────────────────────────────────

def load_document(file_path: str | Path) -> str:
    p = Path(file_path)
    ext = p.suffix.lower()

    if ext == ".pdf":
        # PyMuPDF (fitz) handles RTL/word order far better than pypdf.
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(p))
            parts = [page.get_text("text") or "" for page in doc]
            doc.close()
            return _fix_arabic("\n\n".join(parts))
        except Exception as e:
            logger.error("PDF (PyMuPDF): %s", e)
            # Fallback to pypdf if PyMuPDF fails to open the file.
            try:
                from pypdf import PdfReader
                raw = "\n\n".join(page.extract_text() or "" for page in PdfReader(str(p)).pages)
                return _fix_arabic(raw)
            except Exception as e2:
                logger.error("PDF (pypdf fallback): %s", e2); return ""

    if ext in (".docx", ".doc"):
        try:
            from docx import Document
            text = "\n\n".join(para.text for para in Document(str(p)).paragraphs if para.text.strip())
            return _fix_arabic(text)
        except Exception as e:
            logger.error("DOCX: %s", e); return ""

    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"[Sheet: {sheet}]")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        parts.append(row_text)
            return _fix_arabic("\n".join(parts))
        except Exception as e:
            logger.error("Excel: %s", e); return ""

    if ext == ".csv":
        try:
            import csv, io
            text = p.read_text(encoding="utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(text))
            return _fix_arabic("\n".join(" | ".join(row) for row in reader))
        except Exception as e:
            logger.error("CSV: %s", e); return ""

    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(str(p))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return _fix_arabic("\n\n".join(parts))
        except Exception as e:
            logger.error("PPTX: %s", e); return ""

    if ext in (".html", ".htm"):
        try:
            from html.parser import HTMLParser
            class TX(HTMLParser):
                def __init__(self):
                    super().__init__(); self.text=[]; self.skip=False
                def handle_starttag(self, t, a):
                    if t in ("script","style"): self.skip=True
                def handle_endtag(self, t):
                    if t in ("script","style"): self.skip=False
                def handle_data(self, d):
                    if not self.skip and d.strip(): self.text.append(d.strip())
            tx=TX(); tx.feed(p.read_text(encoding="utf-8",errors="ignore"))
            return _fix_arabic(" ".join(tx.text))
        except Exception as e:
            logger.error("HTML: %s", e); return ""

    if ext == ".json":
        try:
            return json.dumps(json.loads(p.read_text(encoding="utf-8",errors="ignore")), ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error("JSON: %s", e); return ""

    # txt, md — usually already logical order, but normalize just in case.
    return _fix_arabic(p.read_text(encoding="utf-8", errors="ignore"))


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text: return []
    words = text.split()
    chunks, step = [], max(1, chunk_size - overlap)
    for i in range(0, len(words), step):
        c = " ".join(words[i: i + chunk_size])
        if c.strip(): chunks.append(c)
        if i + chunk_size >= len(words): break
    return chunks


# ── FAISS-based Vector Store ──────────────────────────────────────────────────

def build_vector_db(docs_dir, vectordb_path, chunk_size=500, overlap=50) -> dict:
    import faiss
    docs_dir      = Path(docs_dir)
    vectordb_path = Path(vectordb_path)
    vectordb_path.mkdir(parents=True, exist_ok=True)

    embedder  = _get_embedder()
    n_docs    = 0
    all_texts = []
    all_metas = []
    total_size = 0

    for file in docs_dir.iterdir():
        if not file.is_file(): continue
        if file.suffix.lower() not in SUPPORTED_EXTENSIONS:
            logger.warning("Skipping: %s", file.name); continue
        total_size += file.stat().st_size
        text = load_document(file)
        if not text.strip(): logger.warning("Empty: %s", file.name); continue
        chunks = chunk_text(text, chunk_size=chunk_size, overlap=overlap)
        for idx, c in enumerate(chunks):
            all_texts.append(c)
            all_metas.append({"source": file.name, "chunk": idx})
        n_docs += 1
        logger.info("Processed %s: %d chunks", file.name, len(chunks))

    if not all_texts:
        return {"n_documents": 0, "n_chunks": 0, "total_size_mb": 0.0}

    logger.info("Embedding %d chunks...", len(all_texts))
    embeddings = embedder.encode(all_texts, show_progress_bar=False, batch_size=32)
    embeddings = embeddings.astype(np.float32)

    # Normalize for cosine similarity
    faiss.normalize_L2(embeddings)

    # Build FAISS index
    dim   = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)  # Inner product = cosine after normalization
    index.add(embeddings)

    # Save index and metadata
    faiss.write_index(index, str(vectordb_path / "index.faiss"))
    with open(vectordb_path / "metadata.pkl", "wb") as f:
        pickle.dump({"texts": all_texts, "metas": all_metas}, f)

    logger.info("FAISS index built: %d vectors, dim=%d", len(all_texts), dim)
    return {
        "n_documents":   n_docs,
        "n_chunks":      len(all_texts),
        "total_size_mb": round(total_size / (1024 * 1024), 2),
    }


def query_bot(vectordb_path, question, system_prompt="You are a helpful assistant.",
              top_k=3, temperature=0.7, history=None,
              max_tokens=1024, language="auto", fallback_msg=None,
              allow_general=False, similarity_threshold=0.3) -> dict:
    import faiss
    vectordb_path = Path(vectordb_path)
    index_path    = vectordb_path / "index.faiss"
    meta_path     = vectordb_path / "metadata.pkl"

    if not index_path.exists() or not meta_path.exists():
        return {"answer": fallback_msg or "No knowledge base found. Please build the index first.", "sources": []}

    embedder = _get_embedder()
    index    = faiss.read_index(str(index_path))
    with open(meta_path, "rb") as f:
        store = pickle.load(f)
    texts = store["texts"]
    metas = store["metas"]

    # Embed query (repair Arabic in the question too, for consistency).
    q_text = _fix_arabic(question)
    q_emb = embedder.encode([q_text]).astype(np.float32)
    faiss.normalize_L2(q_emb)

    # Search
    k       = min(top_k, len(texts))
    scores, indices = index.search(q_emb, k)
    scores  = scores[0]; indices = indices[0]
    logger.info("Query scores: %s", [round(float(s), 3) for s in scores])

    # Filter by threshold
    results = [(texts[i], metas[i], float(scores[j]))
               for j, i in enumerate(indices)
               if i >= 0 and float(scores[j]) >= similarity_threshold]

    if not results and not allow_general:
        return {"answer": fallback_msg or "لم أجد معلومات كافية في المستندات.", "sources": []}

    if not results:
        results = [(texts[i], metas[i], float(scores[j]))
                   for j, i in enumerate(indices) if i >= 0][:top_k]

    sources = [{"source": m.get("source","?"), "text": t[:300]+"..." if len(t)>300 else t, "score": round(s,3)}
               for t, m, s in results]
    context = "\n\n---\n\n".join(f"[{m.get('source','?')}]\n{t}" for t, m, s in results)

    # Language
    lang_map = {"ar": "Always respond in Arabic.", "en": "Always respond in English."}
    lang_instr = lang_map.get(language, "Respond in the same language as the user.")

    strict = ("Prefer context, but may use general knowledge." if allow_general else
              f"Answer ONLY from context. If not found: '{fallback_msg or 'Not available in documents.'}'")

    messages = [{"role": "system", "content": f"{system_prompt}\n{lang_instr}\n{strict}\n\nCONTEXT:\n{context}"}]
    if history:
        for h in (history or [])[-6:]:
            if h.get("content"):
                role = "assistant" if h.get("role") in ("assistant","model") else "user"
                messages.append({"role": role, "content": h["content"]})
    messages.append({"role": "user", "content": question})

    if not GROQ_KEY:
        logger.error("GROQ_API_KEY is not set")
        return {"answer": fallback_msg or "Connection error. Please try again.", "sources": []}

    try:
        r = httpx.post(GROQ_URL,
            headers={"Authorization": f"Bearer {GROQ_KEY}", "Content-Type": "application/json"},
            json={"model": GROQ_MODEL, "messages": messages, "max_tokens": max_tokens, "temperature": temperature},
            timeout=30)
        if r.status_code == 200:
            return {"answer": r.json()["choices"][0]["message"]["content"], "sources": sources}
        logger.error("Groq %s: %s", r.status_code, r.text[:200])
    except Exception as e:
        logger.error("Groq: %s", e)

    return {"answer": fallback_msg or "Connection error. Please try again.", "sources": []}